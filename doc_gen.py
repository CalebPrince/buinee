"""Generic document rendering Ada calls out to - PDF and PPTX output behind
one small content shape, so any company or vertical can produce a document
without the renderer knowing what business it's for.

Two independent renderers (PDF via reportlab, PPTX via python-pptx), each
fails open rather than crashing the chat: a missing dependency just means
that tool isn't offered, matching outbound_mail.py's is_configured() pattern
for a server-side capability that might not be installed everywhere.

This module only lays out content it's handed - see doc_tools.py for what
decides *what* goes on the page (and where line-item math and tax figures
actually come from, which is never here).
"""

from __future__ import annotations

import io

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable,
    )
    from reportlab.lib.enums import TA_RIGHT
    _REPORTLAB_OK = True
except ImportError:
    _REPORTLAB_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False


def pdf_available() -> bool:
    return _REPORTLAB_OK


def pptx_available() -> bool:
    return _PPTX_OK


NAVY = "0E4F66"
GOLD = "FDB44B"
INK = "1A1A1A"
MUTED = "6A747C"
LIGHT = "F4F7F8"


# --------------------------------------------------------------------- PDF

def build_pdf(spec: dict) -> bytes:
    """spec = {
        "company_name": str, "doc_label": str (e.g. "INVOICE" or "PROPOSAL"),
        "reference": str, "meta_lines": [str, ...] (date, due date, validity...),
        "recipient_lines": [str, ...], "title": str,
        "sections": [
            {"heading": str | None,
             "paragraphs": [str, ...] | None,
             "table": {"headers": [str,...], "rows": [[str,...],...],
                       "emphasize_last_row": bool} | None},
            ...
        ],
        "footer": str,
    }
    Every value is plain text the caller already computed - nothing here
    does arithmetic or invents a number.
    """
    if not _REPORTLAB_OK:
        raise RuntimeError("PDF generation isn't available on this server (reportlab isn't installed).")

    navy, gold, ink, muted, light = (
        colors.HexColor("#" + NAVY), colors.HexColor("#" + GOLD),
        colors.HexColor("#" + INK), colors.HexColor("#" + MUTED), colors.HexColor("#" + LIGHT),
    )

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle("CompanyName", fontName="Helvetica-Bold", fontSize=17, textColor=navy, leading=21))
    styles.add(ParagraphStyle("DocTitle", fontName="Helvetica-Bold", fontSize=20, textColor=ink, leading=24, spaceAfter=4))
    styles.add(ParagraphStyle("Small", fontName="Helvetica", fontSize=9, textColor=muted, leading=13))
    styles.add(ParagraphStyle("SectionHead", fontName="Helvetica-Bold", fontSize=12, textColor=navy, spaceBefore=14, spaceAfter=6))
    styles.add(ParagraphStyle("Body", fontName="Helvetica", fontSize=10.5, textColor=ink, leading=15))
    styles.add(ParagraphStyle("BodyRight", parent=styles["Body"], alignment=TA_RIGHT))

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=22 * mm, rightMargin=22 * mm, topMargin=18 * mm, bottomMargin=18 * mm,
    )
    story = []

    meta_html = f"<b>{spec.get('doc_label', '')}</b>"
    if spec.get("reference"):
        meta_html += f"<br/>No. {spec['reference']}"
    for line in spec.get("meta_lines") or []:
        meta_html += f"<br/>{line}"

    header = Table(
        [[Paragraph(spec.get("company_name", ""), styles["CompanyName"]),
          Paragraph(meta_html, styles["BodyRight"])]],
        colWidths=[100 * mm, 66 * mm],
    )
    header.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP")]))
    story.append(header)
    story.append(Spacer(1, 4))
    story.append(HRFlowable(width="100%", thickness=1.2, color=navy, spaceAfter=14))

    if spec.get("recipient_lines"):
        story.append(Paragraph("<br/>".join(spec["recipient_lines"]), styles["Body"]))
        story.append(Spacer(1, 10))

    if spec.get("title"):
        story.append(Paragraph(spec["title"], styles["DocTitle"]))
        story.append(Spacer(1, 8))

    for section in spec.get("sections") or []:
        if section.get("heading"):
            story.append(Paragraph(section["heading"], styles["SectionHead"]))
        for para in section.get("paragraphs") or []:
            story.append(Paragraph(para, styles["Body"]))
            story.append(Spacer(1, 4))
        table_spec = section.get("table")
        if table_spec and table_spec.get("headers"):
            data = [table_spec["headers"]] + list(table_spec.get("rows") or [])
            ncols = len(table_spec["headers"])
            col_width = 146 / ncols
            tbl = Table(data, colWidths=[col_width * mm] * ncols)
            style = [
                ("BACKGROUND", (0, 0), (-1, 0), navy),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 9.5),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, light]),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#D8DEE2")),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]
            if table_spec.get("emphasize_last_row"):
                style += [
                    ("FONTNAME", (0, -1), (-1, -1), "Helvetica-Bold"),
                    ("BACKGROUND", (0, -1), (-1, -1), gold),
                ]
            tbl.setStyle(TableStyle(style))
            story.append(tbl)
            story.append(Spacer(1, 10))

    if spec.get("footer"):
        story.append(Spacer(1, 10))
        story.append(HRFlowable(width="100%", thickness=0.6, color=muted))
        story.append(Spacer(1, 6))
        story.append(Paragraph(spec["footer"], styles["Small"]))

    doc.build(story)
    return buf.getvalue()


# -------------------------------------------------------------------- PPTX

def build_pptx(spec: dict) -> bytes:
    """spec = {
        "title": str, "subtitle": str, "company_name": str,
        "slides": [
            {"heading": str,
             "bullets": [str, ...] | None,
             "table": {"headers": [...], "rows": [[...]]} | None,
             "stat": {"value": str, "label": str} | None},
            ...
        ],
        "footer": str,
    }
    One shared layout throughout (dark title/closing, light content slides)
    rather than a different template per caller - simple, and reliable
    without a template library to maintain.
    """
    if not _PPTX_OK:
        raise RuntimeError("Presentation generation isn't available on this server (python-pptx isn't installed).")

    navy = RGBColor.from_string(NAVY)
    navy_dark = RGBColor.from_string("0A3A4D")
    gold = RGBColor.from_string(GOLD)
    ink = RGBColor.from_string(INK)
    muted = RGBColor.from_string(MUTED)
    white = RGBColor.from_string("FFFFFF")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide(bg: RGBColor):
        slide = prs.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg
        return slide

    def add_text(slide, text, left, top, width, height, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Arial"
        return box

    # ---- Title slide
    slide = add_slide(navy_dark)
    add_text(slide, (spec.get("title") or "").upper(), 0.8, 2.6, 11.7, 1.0, 34, white, bold=True)
    if spec.get("subtitle"):
        add_text(slide, spec["subtitle"], 0.8, 3.55, 11.7, 0.6, 16, gold)
    if spec.get("company_name"):
        add_text(slide, spec["company_name"], 0.8, 4.15, 11.7, 0.5, 12, RGBColor.from_string("AFC4CE"))

    # ---- Content slides
    for s in spec.get("slides") or []:
        slide = add_slide(white)
        add_text(slide, s.get("heading") or "", 0.6, 0.45, 12, 0.7, 26, navy, bold=True)

        bullets = s.get("bullets")
        table_spec = s.get("table")
        stat = s.get("stat")

        if stat:
            add_text(slide, str(stat.get("value") or ""), 0.6, 1.6, 6, 1.1, 40, navy, bold=True)
            add_text(slide, str(stat.get("label") or ""), 0.6, 2.65, 6, 0.5, 14, muted)
        elif table_spec and table_spec.get("headers"):
            headers = table_spec["headers"]
            rows = table_spec.get("rows") or []
            ncols = len(headers)
            nrows = len(rows) + 1
            gtable = slide.shapes.add_table(nrows, ncols, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.5 * nrows)).table
            for c, h in enumerate(headers):
                cell = gtable.cell(0, c)
                cell.text = str(h)
                cell.text_frame.paragraphs[0].runs[0].font.bold = True
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = white
                cell.fill.solid()
                cell.fill.fore_color.rgb = navy
            for r, row in enumerate(rows, start=1):
                for c, val in enumerate(row):
                    cell = gtable.cell(r, c)
                    cell.text = str(val)
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = ink
        elif bullets:
            box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.3))
            tf = box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = "-  " + item
                run.font.size = Pt(15)
                run.font.color.rgb = ink
                run.font.name = "Arial"
                p.space_after = Pt(10)

    # ---- Closing slide (footer only, same dark treatment as the title)
    if spec.get("footer"):
        slide = add_slide(navy_dark)
        add_text(slide, spec["footer"], 0.8, 3.4, 11.7, 0.8, 14, white)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
